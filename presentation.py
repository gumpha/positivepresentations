# -*- coding: utf-8 -*-
"""
Spyder Editor

This is a temporary script file.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR
from pptx.dml.color import ColorFormat, RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
import random
from random import randrange

def hilo(a, b, c):
    if c < b: b, c = c, b
    if b < a: a, b = b, a
    if c < b: b, c = c, b
    return a + c

def complement(r, g, b):
    k = hilo(r, g, b)
    return tuple(k - u for u in (r, g, b))

def create():
    file = open('positive-words.txt',"r")
    wordList = []
    
    for line in file:
        for word in line.split():
            wordList.append(word)
            
    
    
    prs = Presentation()
    
    for x in range(5):
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        background = slide.background
        fill = background.fill
        fill.solid()
        red = randrange(255)
        green = randrange(255)
        blue = randrange(255)
        fill.fore_color.rgb = RGBColor(red, green, blue)
        width= Inches(10)
        txBox = slide.shapes.add_textbox(Inches(0), Inches(2.5), width, Inches(2.5))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = random.choice(wordList).upper()
        p.font.bold = True
        p.font.size = Pt(120)
        p.alignment = PP_ALIGN.CENTER
        t = complement(red, green, blue) 
        p.font.color.rgb = RGBColor(t[0], t[1], t[2])
    
    
    prs.save('base.pptx')

create()
